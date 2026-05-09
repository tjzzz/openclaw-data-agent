#!/usr/bin/env python3
"""
Data Agent - 环境检查与依赖安装
确保文件读取、PDF处理所需工具就绪。
"""

import os
import sys
import subprocess
import shutil


def check_tool(name, check_cmd=None):
    """检查系统工具是否存在"""
    path = shutil.which(name)
    if path:
        print(f"  ✅ {name}: {path}")
        return True
    # 尝试用 check_cmd 检测
    if check_cmd:
        try:
            result = subprocess.run(check_cmd, shell=True, capture_output=True, text=True)
            if result.returncode == 0:
                print(f"  ✅ {name}: {result.stdout.strip()[:60]}")
                return True
        except Exception:
            pass
    print(f"  ❌ {name}: not found")
    return False


def check_python_module(module_name):
    """检查 Python 模块是否可导入"""
    try:
        __import__(module_name)
        print(f"  ✅ Python module: {module_name}")
        return True
    except ImportError:
        print(f"  ❌ Python module: {module_name}")
        return False


def find_fitz_python():
    """查找能导入 fitz 的 Python 解释器"""
    candidates = [
        sys.executable,
        os.path.expanduser(
            "~/.workbuddy/plugins/marketplaces/cb_teams_marketplace/plugins/document-skills/skills/pdfkit-py/scripts/venv/bin/python3"
        ),
    ]
    for py in candidates:
        if os.path.exists(py):
            try:
                r = subprocess.run([py, "-c", "import fitz; print('ok')"], capture_output=True, text=True, timeout=10)
                if r.returncode == 0:
                    return py
            except Exception:
                pass
    return None


def find_pptx_python():
    """查找能导入 pptx 的 Python 解释器"""
    try:
        import pptx
        return sys.executable
    except ImportError:
        pass
    # 尝试安装
    print("  ⚠️  python-pptx not installed, trying to install...")
    try:
        subprocess.run(
            [sys.executable, "-m", "pip", "install", "python-pptx"],
            capture_output=True, text=True, timeout=120
        )
        import pptx
        return sys.executable
    except Exception:
        return None


def install_pymupdf():
    """安装 PyMuPDF 到系统 Python 或 pdfkit-py venv"""
    print("\n📦 Installing PyMuPDF for PDF reading...")
    
    # 尝试在 pdfkit-py venv 中安装
    pdfkit_venv_pip = os.path.expanduser(
        "~/.workbuddy/plugins/marketplaces/cb_teams_marketplace/plugins/document-skills/skills/pdfkit-py/scripts/venv/bin/pip"
    )
    if os.path.exists(pdfkit_venv_pip):
        result = subprocess.run(
            [pdfkit_venv_pip, "install", "pymupdf"],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            print("  ✅ PyMuPDF installed (pdfkit-py venv)")
            return True

    # Fallback: install to user site
    result = subprocess.run(
        [sys.executable, "-m", "pip", "install", "--user", "pymupdf"],
        capture_output=True, text=True, timeout=120
    )
    if result.returncode == 0:
        print("  ✅ PyMuPDF installed (user site)")
        return True
    
    print("  ❌ Failed to install PyMuPDF")
    return False


def check_file_readers():
    """检查所有文件读取工具"""
    print("\n🔍 File Reader Environment Check")
    print("─" * 40)
    
    doc_ok = check_tool("antiword")
    
    fitz_py = find_fitz_python()
    pdf_ok = fitz_py is not None
    if pdf_ok:
        print(f"  ✅ PDF reader (PyMuPDF): via {fitz_py}")
    else:
        print("  ❌ PDF reader (PyMuPDF): not found")
        print("   尝试安装 PyMuPDF...")
        install_pymupdf()
        fitz_py = find_fitz_python()
        pdf_ok = fitz_py is not None
        if pdf_ok:
            print(f"  ✅ PDF reader (PyMuPDF): installed via {fitz_py}")

    pptx_py = find_pptx_python()
    pptx_ok = pptx_py is not None
    if pptx_ok:
        print(f"  ✅ PPT reader (python-pptx): ok")
    else:
        print("  ❌ PPT reader (python-pptx): not found")
    
    print("─" * 40)
    
    all_ok = doc_ok and pdf_ok and pptx_ok
    if all_ok:
        print("✅ All tools ready")
    else:
        missing = []
        if not doc_ok: missing.append("antiword (.doc)")
        if not pdf_ok: missing.append("PyMuPDF (.pdf)")
        if not pptx_ok: missing.append("python-pptx (.pptx)")
        print(f"⚠️  Missing: {', '.join(missing)}")
    
    return all_ok


def extract_text_from_file(file_path):
    """
    统一文件文本提取接口
    
    Args:
        file_path: 文件路径（.doc, .docx, .pdf, .txt, .md）
    
    Returns:
        str: 提取的文本内容
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.doc':
        # 使用 antiword
        result = subprocess.run(
            ["antiword", "-m", "UTF-8", file_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
        else:
            raise RuntimeError(f"antiword failed: {result.stderr}")

    elif ext == '.pdf':
        # 使用 PyMuPDF（通过能找到 fitz 的 Python 解释器）
        fitz_py = find_fitz_python()
        if not fitz_py:
            raise RuntimeError("PyMuPDF (fitz) not available. Install it first.")
        try:
            result = subprocess.run(
                [fitz_py, "-c", f"""
import fitz, json, sys
doc = fitz.open("{file_path}")
text = ""
for page in doc:
    text += page.get_text()
doc.close()
sys.stdout.write(text)
                """],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                return result.stdout
            else:
                raise RuntimeError(f"PyMuPDF error: {result.stderr}")
        except Exception as e:
            raise RuntimeError(f"PDF read failed: {e}")

    elif ext in ('.docx',):
        # 使用 python-docx 或 antiword
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        except ImportError:
            pass
        # Fallback to antiword
        result = subprocess.run(
            ["antiword", "-m", "UTF-8", file_path],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode == 0:
            return result.stdout
        raise RuntimeError(f"Cannot read .docx: install python-docx or antiword")

    elif ext in ('.txt', '.md'):
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    elif ext in ('.pptx', '.ppt'):
        # 使用 python-pptx
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
            return "\n\n".join(text_parts)
        except ImportError:
            raise RuntimeError("python-pptx not available. Run: python3 scripts/setup_env.py --check")
    
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def main():
    if len(sys.argv) > 1 and sys.argv[1] == "--check":
        ok = check_file_readers()
        sys.exit(0 if ok else 1)
    
    if len(sys.argv) > 1 and sys.argv[1] == "--extract":
        if len(sys.argv) < 3:
            print("Usage: python3 scripts/setup_env.py --extract <file_path>")
            sys.exit(1)
        text = extract_text_from_file(sys.argv[2])
        print(text)
        return

    # Default: run check
    ok = check_file_readers()
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
